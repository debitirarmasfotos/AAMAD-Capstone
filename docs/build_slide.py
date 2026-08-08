"""Generate a native PowerPoint slide of the PMO Program Intelligence Crew
Application Crew and orchestration. Mirrors docs/application-crew-slide.html.

Run: .venv/Scripts/python docs/build_slide.py
Output: docs/application-crew-slide.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# Palette (matches the HTML)
INK      = RGBColor(0x1F, 0x29, 0x33)
MUTED    = RGBColor(0x6B, 0x72, 0x80)
LINE     = RGBColor(0xCB, 0xD5, 0xE1)
SLATE    = RGBColor(0x33, 0x41, 0x55)
INGEST   = RGBColor(0x64, 0x74, 0x8B)
PARALLEL = RGBColor(0x0F, 0x76, 0x6E); PARALLEL_BG = RGBColor(0xE6, 0xF4, 0xF1)
SEQ      = RGBColor(0x1D, 0x4E, 0xD8); SEQ_BG      = RGBColor(0xE6, 0xED, 0xFF)
HITL     = RGBColor(0xB4, 0x53, 0x09); HITL_BG     = RGBColor(0xFD, 0xF1, 0xE0)
FINAL    = RGBColor(0x15, 0x80, 0x3D); FINAL_BG    = RGBColor(0xE7, 0xF4, 0xEC)
DEV      = RGBColor(0x7C, 0x3A, 0xED); DEV_BG      = RGBColor(0xF1, 0xEA, 0xFE)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank


def _set_line(shape, color, width_pt=1.5):
    if color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = color
        shape.line.width = Pt(width_pt)


def box(x, y, w, h, fill, line, radius=0.10):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    _set_line(shp, line)
    shp.shadow.inherit = False
    return shp


def add_text(shape, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=2):
    """runs: list of paragraphs; each paragraph is list of (text, size, bold, color)."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(7); tf.margin_right = Pt(7)
    tf.margin_top = Pt(5); tf.margin_bottom = Pt(5)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        for (text, size, bold, color) in para:
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Segoe UI"
    return shape


def textbox(x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    add_text(tb, runs, align=align, anchor=anchor)
    return tb


def arrow(x, y, w, h):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = SLATE
    _set_line(a, None)
    a.shadow.inherit = False
    try:
        a.adjustments[0] = 0.55  # thinner tail
        a.adjustments[1] = 0.55
    except Exception:
        pass
    return a


def agent_card(x, y, w, h, tag, tag_bg, tag_fg, role, task, border, bg=WHITE, role_color=INK):
    shp = box(x, y, w, h, bg, border)
    tf = shp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top = Pt(7); tf.margin_bottom = Pt(6)
    if tag:
        p0 = tf.paragraphs[0]; p0.space_after = Pt(3)
        r = p0.add_run(); r.text = tag.upper()
        r.font.size = Pt(7.5); r.font.bold = True; r.font.color.rgb = tag_fg
        r.font.name = "Segoe UI"
    p1 = tf.add_paragraph() if tag else tf.paragraphs[0]
    p1.space_after = Pt(3)
    r = p1.add_run(); r.text = role
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = role_color
    r.font.name = "Segoe UI"
    if task:
        p2 = tf.add_paragraph()
        r = p2.add_run(); r.text = task
        r.font.size = Pt(8); r.font.bold = False; r.font.color.rgb = MUTED
        r.font.name = "Segoe UI"
    return shp


# ---- Header ----
textbox(Inches(0.5), Inches(0.28), Inches(12.3), Inches(0.55),
        [[("PMO Program Intelligence Crew", 25, True, INK),
          ("   ·   Application Crew", 25, True, INGEST)]])
textbox(Inches(0.5), Inches(0.82), Inches(12.3), Inches(0.4),
        [[("Runtime: ", 12, False, MUTED), ("claude-agent-sdk", 12, True, SLATE),
          ("    |    End user: ", 12, False, MUTED), ("program manager", 12, True, SLATE),
          ("    |    Orchestration: ", 12, False, MUTED),
          ("hybrid (parallel analysis, sequential synthesis, hierarchical supervisor, HITL gate)", 12, True, SLATE)]])
# header rule
rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.28), Inches(12.33), Pt(2))
rule.fill.solid(); rule.fill.fore_color.rgb = LINE; _set_line(rule, None); rule.shadow.inherit = False

# ---- Flow label ----
textbox(Inches(0.5), Inches(1.42), Inches(12), Inches(0.3),
        [[("RUNTIME FLOW: HOW THE AGENTS USERS INTERACT WITH COORDINATE", 10.5, True, MUTED)]])

# ---- Flow row ----
row_y = Inches(1.95)
card_h = Inches(1.55)
# Ingestion
agent_card(Inches(0.5), Inches(2.25), Inches(1.85), Inches(1.0),
           "Load", None, INGEST, "Ingestion",
           "Reads the project sheet and burn/capacity file into one normalized shared state.",
           INGEST)
arrow(Inches(2.45), Inches(2.65), Inches(0.35), Inches(0.22))

# Parallel stack (3 cards)
px = Inches(2.9)
pw = Inches(3.15)
ph = Inches(0.72)
gap = Inches(0.14)
py0 = Inches(1.72)
parallel_defs = [
    ("Status Rollup", "Derives each priority's RAG from child tasks, citing the rows."),
    ("Capacity / Burn", "Computes utilization and a fit or no-fit signal per workstream."),
    ("Risk & Compliance", "Flags and ranks risks with source evidence and a suggested owner."),
]
for i, (role, task) in enumerate(parallel_defs):
    yy = Emu(int(py0) + i * (int(ph) + int(gap)))
    agent_card(px, yy, pw, ph, None, None, PARALLEL, role, task, PARALLEL, bg=PARALLEL_BG, role_color=PARALLEL)

arrow(Inches(6.15), Inches(2.65), Inches(0.35), Inches(0.22))

# Narrative
agent_card(Inches(6.6), Inches(2.25), Inches(1.95), Inches(1.0),
           "Synthesize", None, SEQ, "Narrative",
           "Builds the DRAFT executive readout from all analysis outputs.",
           SEQ, bg=SEQ_BG, role_color=SEQ)
arrow(Inches(8.65), Inches(2.65), Inches(0.35), Inches(0.22))

# HITL
agent_card(Inches(9.1), Inches(2.25), Inches(1.95), Inches(1.0),
           "Human", None, HITL, "HITL Gate",
           "Program manager approves, edits, or rejects before anything is final.",
           HITL, bg=HITL_BG, role_color=HITL)
arrow(Inches(11.15), Inches(2.65), Inches(0.35), Inches(0.22))

# FINAL
final = box(Inches(11.6), Inches(2.4), Inches(1.25), Inches(0.7), FINAL_BG, FINAL)
add_text(final, [[("FINAL", 14, True, FINAL)], [("Approved readout", 8, False, MUTED)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ---- Supervisor band ----
sup = box(Inches(0.5), Inches(3.75), Inches(12.33), Inches(0.7), SLATE, None, radius=0.12)
tf = sup.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Pt(12)
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Supervisor / Orchestrator   "
r.font.size = Pt(12.5); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Segoe UI"
r = p.add_run(); r.text = "holds shared state across all agents, routes work, and enforces the approval gate.   "
r.font.size = Pt(11.5); r.font.color.rgb = WHITE; r.font.name = "Segoe UI"
r = p.add_run(); r.text = "(hierarchical, spans the whole flow)"
r.font.size = Pt(10.5); r.font.italic = True; r.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1); r.font.name = "Segoe UI"

# ---- Footer: Development Crew ----
dev = box(Inches(0.5), Inches(4.75), Inches(8.4), Inches(2.3), DEV_BG, DEV, radius=0.06)
dev.line.dash_style = None
tf = dev.text_frame; tf.word_wrap = True
tf.margin_left = Pt(12); tf.margin_right = Pt(12); tf.margin_top = Pt(10)
p = tf.paragraphs[0]; p.space_after = Pt(8)
r = p.add_run(); r.text = "DEVELOPMENT CREW  ·  BUILDS THE SYSTEM (IDE / BUILD TIME, NOT SHIPPED)"
r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = DEV; r.font.name = "Segoe UI"
for label, people in [
    ("Define: ", "product-mgr (MRD, PRD), system-arch (SAD)"),
    ("Build: ", "project-mgr, backend-eng, frontend-eng, integration-eng"),
    ("Validate / Deliver: ", "qa-eng, security-eng, devops-eng"),
]:
    p = tf.add_paragraph(); p.space_after = Pt(5)
    r = p.add_run(); r.text = label; r.font.size = Pt(11.5); r.font.bold = True; r.font.color.rgb = DEV; r.font.name = "Segoe UI"
    r = p.add_run(); r.text = people; r.font.size = Pt(11.5); r.font.color.rgb = INK; r.font.name = "Segoe UI"
p = tf.add_paragraph(); p.space_before = Pt(6)
r = p.add_run(); r.text = "The Development Crew builds the thing. The Application Crew above is the thing."
r.font.size = Pt(10.5); r.font.italic = True; r.font.color.rgb = MUTED; r.font.name = "Segoe UI"

# ---- Footer: Legend ----
leg = box(Inches(9.05), Inches(4.75), Inches(3.78), Inches(2.3), WHITE, LINE, radius=0.06)
tf = leg.text_frame; tf.word_wrap = True
tf.margin_left = Pt(12); tf.margin_right = Pt(10); tf.margin_top = Pt(10)
p = tf.paragraphs[0]; p.space_after = Pt(9)
r = p.add_run(); r.text = "PATTERN LEGEND"; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = SLATE; r.font.name = "Segoe UI"
legend_items = [
    (PARALLEL, "Parallel: independent analysis runs concurrently"),
    (SEQ, "Sequential: Narrative depends on all analysis"),
    (HITL, "Human-in-the-loop: hard stop before final"),
    (SLATE, "Hierarchical: Supervisor owns shared state"),
]
for color, text in legend_items:
    p = tf.add_paragraph(); p.space_after = Pt(7)
    r = p.add_run(); r.text = "■  "; r.font.size = Pt(11); r.font.color.rgb = color; r.font.name = "Segoe UI"
    r = p.add_run(); r.text = text; r.font.size = Pt(10.5); r.font.color.rgb = INK; r.font.name = "Segoe UI"

out = "docs/application-crew-slide.pptx"
prs.save(out)
print("saved", out)
